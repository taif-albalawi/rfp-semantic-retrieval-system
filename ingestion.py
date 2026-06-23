from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import uuid
import re
import json

from chunking import chunk_text

data_folder = Path("data")

def clean_text(text):
    text = re.sub(r"\s+", " ", text)
    text = text.replace("\x00", "")
    return text.strip()


def extract_text(file_path):

    # PDF
    if file_path.suffix == ".pdf":
        reader = PdfReader(file_path)
        text = ""

        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text()

        return text

    # DOCX
    elif file_path.suffix == ".docx":
        doc = Document(file_path)

        text = ""

        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"

        return text

    # XLSX
    elif file_path.suffix == ".xlsx":

        sheets = pd.read_excel(file_path, sheet_name=None)

        text = ""

        for sheet_name, df in sheets.items():
            text += df.to_string()

        return text
    
    # PPTX
    elif file_path.suffix == ".pptx":
        presentation = Presentation(file_path)

        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    return ""


supported_extensions = [".pdf", ".docx", ".xlsx", ".pptx"]

files = [
    f for f in data_folder.rglob("*")
    if f.is_file() and f.suffix.lower() in supported_extensions
]

seen_files = set()

all_chunks = []

for file in files:

    file_key = (file.name, file.stat().st_size)

    if file_key in seen_files:
        continue

    seen_files.add(file_key)

    text = extract_text(file)
    text = clean_text(text)

    if not text:
        continue

    doc_id = str(uuid.uuid4())

    chunks = chunk_text(text)

    for c in chunks:
        all_chunks.append({
        "doc_id": doc_id,
        "source": file.name,
        "file_type": file.suffix.lower(),
        "folder": file.parent.name,
        "chunk_id": c["chunk_id"],
        "content": c["content"]
        })

    print(f"\nFILE: {file.name}")
    print(f"TOTAL CHUNKS: {len(chunks)}")
    print("FIRST CHUNK:\n", chunks[0]["content"][:200])

with open("chunks.json", "w", encoding="utf-8") as f:
    json.dump(all_chunks, f, ensure_ascii=False, indent=2)

print("\nSaved chunks to chunks.json")
